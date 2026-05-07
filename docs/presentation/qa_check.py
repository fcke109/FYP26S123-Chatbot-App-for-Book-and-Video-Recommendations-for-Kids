"""Layout sanity checker for the deck — flags potential overlaps and overflows.

Slide is 13.3" wide x 7.5" tall (LAYOUT_WIDE).
"""
from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400
SLIDE_W = 13.3
SLIDE_H = 7.5

prs = Presentation("LittleDino_CF_Defense.pptx")

issues = []
for idx, slide in enumerate(prs.slides, 1):
    rects = []
    for shape in slide.shapes:
        if shape.left is None:
            continue
        x = shape.left / EMU_PER_INCH
        y = shape.top / EMU_PER_INCH
        w = shape.width / EMU_PER_INCH
        h = shape.height / EMU_PER_INCH

        # Off-slide?
        if x < -0.05 or y < -0.05 or x + w > SLIDE_W + 0.05 or y + h > SLIDE_H + 0.05:
            issues.append(
                f"Slide {idx}: shape '{shape.shape_type}' off-slide "
                f"(x={x:.2f}, y={y:.2f}, w={w:.2f}, h={h:.2f})"
            )

        # Margin from edges (warn if < 0.4")
        if shape.has_text_frame and shape.text_frame.text.strip():
            if x < 0.3 or x + w > SLIDE_W - 0.3:
                # only warn if it's a body content shape (not the accent bar at x=0)
                if w > 1.0:
                    issues.append(
                        f"Slide {idx}: text near edge "
                        f"(x={x:.2f}, w={w:.2f}): \"{shape.text_frame.text[:50]}...\""
                    )

        rects.append((shape, x, y, w, h))

    # Overlap detection — pairwise (only flag if BOTH have text)
    for i in range(len(rects)):
        for j in range(i + 1, len(rects)):
            sa, xa, ya, wa, ha = rects[i]
            sb, xb, yb, wb, hb = rects[j]
            if not (sa.has_text_frame and sb.has_text_frame):
                continue
            if not (sa.text_frame.text.strip() and sb.text_frame.text.strip()):
                continue
            # check intersection
            if xa < xb + wb and xa + wa > xb and ya < yb + hb and ya + ha > yb:
                # Overlap — but only flag if SIGNIFICANT (>30% of smaller box)
                ix = max(0, min(xa + wa, xb + wb) - max(xa, xb))
                iy = max(0, min(ya + ha, yb + hb) - max(ya, yb))
                area = ix * iy
                smaller = min(wa * ha, wb * hb)
                if area > 0.3 * smaller:
                    issues.append(
                        f"Slide {idx}: significant text overlap "
                        f"\"{sa.text_frame.text[:30]}\" vs \"{sb.text_frame.text[:30]}\""
                    )

print(f"Inspected {len(prs.slides)} slides")
print(f"Found {len(issues)} potential issues:")
for issue in issues:
    print(f"  - {issue}")
